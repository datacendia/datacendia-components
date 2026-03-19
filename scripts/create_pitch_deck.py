"""
Datacendia Investor Pitch Deck Generator
For meeting with Alejandro Cremades - January 12, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Helper to create RGB color
def rgb_color(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Colors
DARK_BG = rgb_color(10, 10, 10)
GOLD = rgb_color(212, 175, 55)
WHITE = rgb_color(255, 255, 255)
LIGHT_GRAY = rgb_color(180, 180, 180)
DARK_GRAY = rgb_color(40, 40, 40)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = GOLD
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.5))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = LIGHT_GRAY
    
    # Bullets
    start_y = 1.8 if subtitle else 1.4
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(start_y), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
    ltf = left_title_box.text_frame
    lp = ltf.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(24)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    
    # Left bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.8), Inches(5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = ltf.paragraphs[0]
        else:
            p = ltf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
    rtf = right_title_box.text_frame
    rp = rtf.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(24)
    rp.font.bold = True
    rp.font.color.rgb = WHITE
    
    # Right bullets
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = rtf.paragraphs[0]
        else:
            p = rtf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Metrics boxes
    num_metrics = len(metrics)
    box_width = 2.5
    total_width = num_metrics * box_width + (num_metrics - 1) * 0.3
    start_x = (13.33 - total_width) / 2
    
    for i, (number, label, sublabel) in enumerate(metrics):
        x = start_x + i * (box_width + 0.3)
        
        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(box_width), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_GRAY
        box.line.color.rgb = GOLD
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(box_width), Inches(1))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = str(number)
        np.font.size = Pt(48)
        np.font.bold = True
        np.font.color.rgb = GOLD
        np.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.6), Inches(box_width), Inches(0.6))
        ltf = label_box.text_frame
        lp = ltf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(16)
        lp.font.bold = True
        lp.font.color.rgb = WHITE
        lp.alignment = PP_ALIGN.CENTER
        
        # Sublabel
        if sublabel:
            sub_box = slide.shapes.add_textbox(Inches(x), Inches(4.2), Inches(box_width), Inches(1))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = sublabel
            sp.font.size = Pt(12)
            sp.font.color.rgb = LIGHT_GRAY
            sp.alignment = PP_ALIGN.CENTER
    
    return slide

def create_pitch_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ==========================================================================
    # SLIDE 1: Title
    # ==========================================================================
    add_title_slide(prs, "DATACENDIA", "Sovereign Enterprise Intelligence Platform")
    
    # ==========================================================================
    # SLIDE 2: The Problem
    # ==========================================================================
    add_content_slide(prs, "The Problem", [
        "Enterprise AI is a black box — regulators can't audit it, boards can't explain it",
        "Cloud AI vendors have custody of your data (CLOUD Act exposure, vendor lock-in)",
        "Existing BI shows what happened — not what to do about it",
        "No audit trail = no accountability = no trust",
        "Banks, defense, healthcare, legal firms CANNOT use cloud AI for sensitive decisions"
    ], "Regulated industries are locked out of AI")
    
    # ==========================================================================
    # SLIDE 3: The Solution
    # ==========================================================================
    add_content_slide(prs, "The Solution: Sovereign Intelligence", [
        "AI that runs on YOUR infrastructure (on-prem, private cloud, air-gapped)",
        "Encryption keys stay in YOUR KMS/HSM — we never see your data",
        "Every decision logged to an immutable ledger YOU control",
        "Cryptographically signed decision packets for regulators",
        "Multi-agent deliberation: 45 specialized AI agents debate like a boardroom"
    ], "Not BI. Not another AI API. A decision layer that proves why it decided what it decided.")
    
    # ==========================================================================
    # SLIDE 4: Product Metrics
    # ==========================================================================
    add_metrics_slide(prs, "Platform Metrics", [
        ("1,464", "Automated Tests", "100% pass rate"),
        ("45", "Council Agents", "CEO, CFO, CISO, Legal..."),
        ("16", "Security Tests", "OWASP LLM Top 10"),
        ("26", "Languages", "Full i18n"),
        ("4", "Deploy Modes", "Cloud/Private/On-Prem/Air-Gap")
    ])
    
    # ==========================================================================
    # SLIDE 5: Core Products
    # ==========================================================================
    add_two_column_slide(prs, "Core Products",
        "The Council™",
        [
            "14+ specialized agents deliberate on decisions",
            "CFO, CISO, Risk, Strategy, Legal, Red Team",
            "Not a chatbot — a boardroom debate",
            "Structured dissent & minority opinions logged",
            "Cryptographic signatures on every decision"
        ],
        "CendiaChronos™",
        [
            "Enterprise Time Machine",
            "Simulate divergent futures before committing",
            "What-if scenario modeling",
            "Pivotal moment detection",
            "Timeline visualization"
        ]
    )
    
    # ==========================================================================
    # SLIDE 6: Trust Layer
    # ==========================================================================
    add_content_slide(prs, "The Trust Layer", [
        "Decision DNA™ — Immutable audit trail with cryptographic proof",
        "Regulator's Receipt — One-click evidence packets for auditors",
        "CendiaCrucible™ — Adversarial stress testing (red-team your own AI)",
        "CendiaOversight™ — Compliance monitoring (SOC 2, HIPAA, NIST 800-53)",
        "CendiaDissent™ — Formal dissent filing with retaliation protection"
    ], "When the auditor asks 'Why?', hand them this.")
    
    # ==========================================================================
    # SLIDE 7: Four Verticals
    # ==========================================================================
    add_two_column_slide(prs, "Four Target Verticals",
        "Legal & Law Firms",
        [
            "12-step legal decision workflow",
            "Matter Lead, Research Counsel, Litigation Strategist",
            "Privilege & confidentiality protection",
            "ABA Model Rules compliance",
            "Regulator's Receipt for court review"
        ],
        "Healthcare & Pharma",
        [
            "FDA 21 CFR Part 11 compliant",
            "Clinical trial decision support",
            "HIPAA-aligned architecture",
            "Genomics integration ready",
            "IRB-ready evidence packets"
        ]
    )
    
    # ==========================================================================
    # SLIDE 8: More Verticals
    # ==========================================================================
    add_two_column_slide(prs, "Four Target Verticals (cont.)",
        "Financial Services",
        [
            "SOC 2 aligned architecture",
            "Investment committee workflows",
            "Risk assessment with dissent tracking",
            "Regulatory compliance (SEC, FINRA)",
            "Zero-copy data architecture"
        ],
        "Government & Defense",
        [
            "Air-gapped deployment",
            "NIST 800-53 controls",
            "FedRAMP-ready architecture",
            "TPM hardware attestation",
            "Classified environment support"
        ]
    )
    
    # ==========================================================================
    # SLIDE 9: Sovereign Architecture
    # ==========================================================================
    add_content_slide(prs, "11 Sovereign Architecture Patterns", [
        "Data Diode — Unidirectional data ingest with quarantine",
        "Local RLHF — Zero-cloud reinforcement learning",
        "QR Air-Gap Bridge — Animated QR for zero-media transfer",
        "TPM Attestation — Hardware-signed decisions",
        "Time-Lock Crypto — Embargoed decisions with cryptographic release",
        "Federated Mesh — Multi-site learning via sneakernet",
        "Portable USB Instance — Bootable deployment"
    ], "Enterprise-grade patterns for the most demanding environments")
    
    # ==========================================================================
    # SLIDE 10: Market Opportunity
    # ==========================================================================
    add_content_slide(prs, "Market Opportunity", [
        "Enterprise AI software market: $297.9B by 2027 (Gartner, 19.1% CAGR)",
        "Government sector: Largest AI spend at $70B+ by 2027 (Gartner)",
        "85% of executives say compliance has become more complex in 3 years (PwC 2025)",
        "Gap: No sovereign, auditable, multi-agent decision platform exists",
        "First-mover advantage in 'Explainable Sovereign AI' category"
    ], "Regulated industries are desperate for AI they can actually use")
    
    # ==========================================================================
    # SLIDE 11: Business Model
    # ==========================================================================
    add_content_slide(prs, "Business Model", [
        "Sandbox Wargame: $5,000 (2-week technical evaluation)",
        "Pilot Program: $50,000 (90-day proof of value)",
        "Department License: $180,000/year + $25K-$75K onboarding",
        "Enterprise License: Custom annual licensing (multi-department)",
        "Sovereign Tier: Classified pricing (air-gapped, high-assurance)",
        "Land & Expand: Pilot → Department → Enterprise → Sovereign"
    ], "Pricing from datacendia.com/pricing.html")
    
    # ==========================================================================
    # SLIDE 12: Competitive Advantage
    # ==========================================================================
    add_content_slide(prs, "Why We Win", [
        "Only platform with TRUE sovereignty (your infra, your keys, your ledger)",
        "Only platform with multi-agent deliberation (not single-model chatbot)",
        "Only platform with cryptographic audit trail (Regulator's Receipt)",
        "Only platform with formal dissent tracking (minority opinions matter)",
        "Built by enterprise architects, not consumer AI researchers"
    ], "See next slides for detailed competitor comparisons")
    
    # ==========================================================================
    # SLIDE 12b: Datacendia vs Harvey AI
    # ==========================================================================
    add_content_slide(prs, "Datacendia vs Harvey AI", [
        "Harvey AI: $8B valuation, $760M raised in 2025, 50 AmLaw 100 firms",
        "",
        "Harvey: Legal ONLY | Datacendia: Legal + Healthcare + Finance + Gov",
        "Harvey: Cloud ONLY | Datacendia: Cloud + Private + On-Prem + Air-Gap",
        "Harvey: Single model | Datacendia: 45-agent multi-agent deliberation",
        "Harvey: No audit trail | Datacendia: Regulator's Receipt (cryptographic)",
        "Harvey: No dissent | Datacendia: CendiaDissent™ formal tracking"
    ], "Harvey proves the market. Datacendia serves what regulated enterprises need.")
    
    # ==========================================================================
    # SLIDE 12c: Datacendia vs Palantir (Functionality)
    # ==========================================================================
    add_content_slide(prs, "Datacendia vs Palantir: Functionality", [
        "PALANTIR STRENGTHS (we acknowledge):",
        "   • Foundry: Best-in-class data integration & ontology modeling",
        "   • Gotham: Proven defense/intel platform (20+ years)",
        "   • AIP: LLM workflows, Agent Studio, massive scale",
        "",
        "DATACENDIA DIFFERENTIATORS:",
        "   • Multi-agent deliberation (45 agents debate, not single LLM)",
        "   • Regulator's Receipt (cryptographic audit trail for compliance)",
        "   • CendiaDissent™ (formal minority opinion tracking)",
        "   • 90-day pilot vs 6-12 month Palantir deployment",
        "   • $50K entry vs $1M+ Palantir minimum"
    ], "Palantir = data platform. Datacendia = decision intelligence + audit proof.")
    
    # ==========================================================================
    # SLIDE 12d: Market Positioning
    # ==========================================================================
    add_content_slide(prs, "Market Positioning", [
        "PALANTIR: $100M+ enterprises, government, defense",
        "   → $423B market cap, $1M+ deals, 6-12 month sales cycles",
        "",
        "HARVEY AI: AmLaw 100 law firms, legal-only",
        "   → $8B valuation, cloud-only, single vertical",
        "",
        "DATACENDIA: Mid-market regulated enterprises ($10M-$500M)",
        "   → $50K pilots, 4 verticals, sovereign deployment options",
        "   → The 'Palantir for the rest of us' with audit-proof decisions"
    ], "We don't compete with Palantir. We serve who they can't reach.")
    
    # ==========================================================================
    # SLIDE 13: Team
    # ==========================================================================
    add_content_slide(prs, "Team", [
        "Stuart Rainey — Founder & Architect",
        "   • 8 years as Senior Data Analyst",
        "   • Masters Degree in Big Data Analytics",
        "   • Deep expertise in enterprise data systems and AI governance",
        "",
        "Seeking: CTO, VP Sales, VP Customer Success"
    ])
    
    # ==========================================================================
    # SLIDE 14: Traction & Milestones
    # ==========================================================================
    add_content_slide(prs, "Current Status & Roadmap", [
        "✓ DONE: Full platform built (1,464 tests, 45 agents, 4 deploy modes)",
        "✓ DONE: Live demo at datacendia.com",
        "→ Q1 2026: First 3 pilot customers (Legal, Healthcare, Finance)",
        "→ Q2 2026: SOC 2 Type II certification",
        "→ Q3 2026: First Fortune 500 enterprise customer",
        "→ Q4 2026: $1M ARR milestone"
    ])
    
    # ==========================================================================
    # SLIDE 15: The Ask
    # ==========================================================================
    add_title_slide(prs, "The Ask", "Pre-Seed: $1.5M at $7M Pre-Money (17.6% Dilution)")
    
    # ==========================================================================
    # SLIDE 16: Use of Funds
    # ==========================================================================
    add_content_slide(prs, "Use of Funds ($1.5M, 18-month runway)", [
        "Engineering (35% / $525K) — 2 senior engineers",
        "Sales & Marketing (25% / $375K) — 1 enterprise sales + marketing",
        "Compliance & Certs (15% / $225K) — SOC 2 Type II, ISO 27001",
        "Infrastructure (10% / $150K) — Cloud, hosting, dev tools",
        "Legal & Admin (5% / $75K) — Legal, accounting, insurance",
        "Founder Salary (10% / $150K) — $5,000/month × 18 months"
    ])
    
    # ==========================================================================
    # SLIDE 17: References
    # ==========================================================================
    add_content_slide(prs, "References", [
        "1. Gartner (Aug 2025). 'AI Agent Enterprise Forecast.'",
        "   → 40% of enterprise apps will embed AI agents by 2026 (up from <5%)",
        "",
        "2. IBM (2025). 'Cost of a Data Breach Report.'",
        "   → 97% of AI breaches lacked access controls; 63% lacked governance",
        "",
        "3. PRNewswire/Iron Software (Sept 2025). 'Data Sovereignty Revolution.'",
        "   → 300% surge in enterprise demand for air-gapped AI solutions",
        "",
        "4. CRN/Xpert Digital (2025). 'Cloud Outage Costs.'",
        "   → Azure ($16B), AWS ($581M), Cloudflare ($250M+) — 2025 alone",
        "",
        "5. Clio (2024). 'Law Firm Data Breach Report.'",
        "   → Law firm breach cost: $5.08M avg (+10% YoY)"
    ])
    
    # ==========================================================================
    # SLIDE 18: Contact
    # ==========================================================================
    slide = add_title_slide(prs, "Let's Talk", "contact@datacendia.com")
    
    # Add website
    web_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.33), Inches(0.5))
    tf = web_box.text_frame
    p = tf.paragraphs[0]
    p.text = "datacendia.com"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), '..', 'Datacendia_Pitch_Deck_Alejandro_Cremades.pptx')
    prs.save(output_path)
    print(f"Pitch deck saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_pitch_deck()
